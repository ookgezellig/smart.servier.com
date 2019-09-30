## Environment to run this script in
# Windows 10 Pro machine
# MS Office 2016 installed on machine
# Python 3.7

import lxml
import os, os.path, glob
import json
from pprint import pprint
from lxml import html
import requests
from pptx import Presentation
import win32com.client
from comtypes.client import Constants, CreateObject



##### STEP 1: Extract texts and metadata from every ppt file and save them in 1 json file
def extract_texts_metadata(files,dir,filename):
    dict2 = {}
    dict = {}
    for file in files:
        pptname = file.split("\\")[-1] #Strip the file path, get only the filename
        dict2[pptname]={}
        #Extract texts from slides
        #Codeblock taken from https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
        prs = Presentation(file)
        # text_runs will be populated with a list of strings, one for each text run in presentation
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        #Build dict
        dict["filename"] = pptname
        dict["slidetext"] = text_runs
        #Extract metadata from ppt files
        dict["metadata-subject"] = prs.core_properties.subject
        dict["metadata-comments"] = prs.core_properties.comments
        dict2[pptname].update(dict)
    #Write dict to json file
    output=os.path.join(dir, filename)
    with open(output, 'w') as jsonfile:
        json.dump(dict2, jsonfile)

### STEP 2 Manually correct the json file
# TODO: Manual step to correct json file by hand to get rid of errors

#### STEP 3: After manually correction, use corrected json file and os.rename()
# to rename ppts with descriptive filenames and move them to 'ppt' folder

def rename_ppt(files,dir,filename):
    from os.path import dirname as up
    from shutil import copyfile
    # Open corrected dict
    input = os.path.join(dir, filename)
    with open(input) as jsonfile:
        dict = json.load(jsonfile)
    for file in files:
        pptname=file.split("\\")[-1]
        # Build new, descriptive filenames for ppt
        pptnewname = ""
        pptnewname += dict[pptname]['metadata-subject'] + " - "
        if dict[pptname]['metadata-comments'] != dict[pptname]['metadata-subject'] :
            pptnewname += dict[pptname]['metadata-comments'] + " - "
        if dict[pptname]['slidetext'][0] != dict[pptname]['metadata-comments'] :
            pptnewname += dict[pptname]['slidetext'][0]
        pptnewname2 = pptnewname.strip(" - ")
        pptnewname2 +='.pptx'
        # Move to 'ppt' folder
        two_up = up(up(file))  # Go to grandparent folder, see https://stackover flow.com/questions/27844088/python-get-directory-two-levels-up
        movedfile = two_up + "\\ppt\\" + pptnewname2
        print(file + " --> " + movedfile)
        copyfile(file,movedfile)

### STEP $4====
# Convert newly named ppts to jpgs
#https://stackoverflow.com/questions/52258446/using-file-format-constants-when-saving-powerpoint-presentation-with-comtypes
#https://docs.microsoft.com/en-us/previous-versions/office/developer/office-2010/ff760025(v=office.14)
#https://docs.microsoft.com/en-us/office/vba/api/word.wdsaveformat

def convert_ppt2images(files, dir):
    from os.path import dirname as up
    import shutil
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    powerpoint2 = CreateObject("Powerpoint.Application")
    pp_constants = Constants(powerpoint2)
    powerpoint.Visible = 1
    for file in files:
        print('file= '+ file)
        filepath = os.path.splitext(file)[0]
        print('filepath= '+ filepath)
        newname = filepath + ".jpg"
        deck = powerpoint.Presentations.Open(file)
        deck.SaveAs(newname, pp_constants.ppSaveAsJPG)
        deck.Close()

        # Output is a folder with 1 fixed-name jpg file (media/Bones/Bone_fractures/ppt/Bone fractures - Ankle fractures/Dia1.JPG)
        # We need to
        # 1) copy this folder (and its content) to the '../jpg-png' folder
        # 2) there rename Dia1.JPG to 'Bone fractures - Ankle fractures.jpg'
        # 3) remove the folder media/Bones/Bone_fractures/ppt/Bone fractures - Ankle fractures/ + its content
        # We can improve the variable names : path, newname, newdir, oldfile, newfile --- rather confusing I say old chap!
        # But hey, it works!

        lastdir_in_filepath= filepath.split("\\")[-1]
        print('lastdir_in_filepath= '+lastdir_in_filepath)

        newdir= dir + "\\" + lastdir_in_filepath
        print('newdir= '+ newdir)

        if not os.path.exists(newdir):
            os.mkdir(newdir)

        oldfile=os.path.join(filepath, 'Dia1.JPG')
        print('oldfile= '+ oldfile)
        newfile=os.path.join(newdir, lastdir_in_filepath + '.jpg')
        print('newfile= '+ newfile)
        print('---------------------')
        if not os.path.exists(newfile):
            os.rename(oldfile,newfile)

        shutil.rmtree(filepath)
    powerpoint.Quit()

#### Main stuff #####

#config paths/ urls
currentdir = os.path.dirname(os.path.realpath(__file__)) # Path of this .py file
basedir= "Bones" # Change this according to the targeted base folder
subdir = "Skeleton and bones" # Change this according to the targeted subfolder
homedir= currentdir + "\\" + basedir + "\\" + subdir + "\\"

pptsourcedir= homedir + "\\ppt-source\\" # Folder with (single slide) ppts with the original, non-descriptive file names
files = glob.glob(os.path.join(pptsourcedir, '*.pptx'))
#extract_texts_metadata(files,homedir,'folder-content.json')

pptdir = homedir + "\\ppt\\" #Folder with ppts with improved, descriptive file names
#if not os.path.exists(pptdir):
#    os.makedirs(pptdir)
#rename_ppt(files,homedir,'folder-content-corrected.json')

files2 = glob.glob(os.path.join(pptdir, '*.pptx'))

imagesdir = homedir + "\\images\\" #Folder with jpg (main slide, English filename) & pngs (cut-outs, French filenames) with improved, descriptive file names
#if not os.path.exists(imagesdir):
#    os.makedirs(imagesdir)
#convert_ppt2images(files2, imagesdir)



